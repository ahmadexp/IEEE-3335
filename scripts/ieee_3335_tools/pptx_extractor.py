"""PPTX extraction and conversion utilities."""

from pathlib import Path
from typing import Optional
from rich.console import Console

console = Console()


def extract_pptx_to_markdown(
    pptx_path: Path, 
    output_path: Optional[Path] = None,
    extract_images: bool = True
) -> Path:
    """
    Extract PPTX to markdown using python-pptx.
    
    Args:
        pptx_path: Path to input PPTX file
        output_path: Optional output path (default: {input_name}.md)
        extract_images: Whether to extract and save images
        
    Returns:
        Path to the created markdown file
    """
    try:
        from pptx import Presentation
    except ImportError as e:
        raise ImportError(
            f"Required library not installed: {e}\nPlease install with: uv add python-pptx pillow"
        )
    
    if output_path is None:
        output_path = pptx_path.with_suffix('.md')
    
    console.print(f"[bold]Converting PPTX to Markdown:[/bold] {pptx_path}")
    
    prs = Presentation(str(pptx_path))
    markdown_lines = []
    
    # Create images directory if extracting images
    images_dir = None
    if extract_images:
        images_dir = output_path.parent / f"{output_path.stem}_images"
        images_dir.mkdir(exist_ok=True)
    
    image_counter = 0
    
    for slide_num, slide in enumerate(prs.slides, 1):
        markdown_lines.append(f"# Slide {slide_num}")
        markdown_lines.append("")
        
        # Extract text from shapes
        for shape in slide.shapes:
            # Handle text
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                
                # Try to detect if it's a title
                if hasattr(shape, "placeholder_format"):
                    if shape.placeholder_format.type in [1, 3]:  # Title or Center Title
                        markdown_lines.append(f"## {text}")
                    else:
                        markdown_lines.append(text)
                else:
                    markdown_lines.append(text)
                markdown_lines.append("")
            
            # Handle images
            if extract_images and shape.shape_type == 13:  # Picture
                try:
                    image = shape.image
                    image_bytes = image.blob
                    
                    # Determine extension
                    ext = image.ext
                    image_counter += 1
                    image_filename = f"slide_{slide_num}_image_{image_counter}.{ext}"
                    image_path = images_dir / image_filename
                    
                    # Save image
                    image_path.write_bytes(image_bytes)
                    
                    # Add markdown reference
                    relative_path = f"{output_path.stem}_images/{image_filename}"
                    markdown_lines.append(f"![Image {image_counter}]({relative_path})")
                    markdown_lines.append("")
                except Exception as e:
                    console.print(f"[yellow]Warning: Could not extract image from slide {slide_num}: {e}[/yellow]")
            
            # Handle tables
            if shape.has_table:
                table = shape.table
                markdown_lines.append("")
                for i, row in enumerate(table.rows):
                    cells = [cell.text.strip() for cell in row.cells]
                    markdown_lines.append("| " + " | ".join(cells) + " |")
                    if i == 0:
                        markdown_lines.append("| " + " | ".join(["---"] * len(cells)) + " |")
                markdown_lines.append("")
        
        # Add separator between slides
        markdown_lines.append("---")
        markdown_lines.append("")
    
    # Write output
    markdown_content = "\n".join(markdown_lines)
    output_path.write_text(markdown_content, encoding='utf-8')
    
    success_msg = f"[green]✓ Successfully converted to markdown:[/green] {output_path}"
    if extract_images and image_counter > 0:
        success_msg += f"\n[green]✓ Extracted {image_counter} image(s) to:[/green] {images_dir}"
    console.print(success_msg)
    
    return output_path
